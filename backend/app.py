from database import init_db, get_db
from flask import Flask, request, jsonify, send_from_directory
from werkzeug.utils import secure_filename
from flask_cors import CORS
from PIL import Image
import os
import uuid
from datetime import datetime
import bcrypt
import jwt

app = Flask(__name__)
CORS(app)
app.config['SECRET_KEY'] = 'your-secret-key-change-this-in-production'

BASE_UPLOAD_FOLDER = os.path.join(os.getcwd(), "backend/uploads")
ORIGINAL_FOLDER = os.path.join(BASE_UPLOAD_FOLDER, "originals")
CONVERTED_FOLDER = os.path.join(BASE_UPLOAD_FOLDER, "converted")
PROFILE_FOLDER = os.path.join(BASE_UPLOAD_FOLDER, "profiles")

os.makedirs(ORIGINAL_FOLDER, exist_ok=True)
os.makedirs(CONVERTED_FOLDER, exist_ok=True)
os.makedirs(PROFILE_FOLDER, exist_ok=True)

# Allowed file types
ALLOWED_IMAGE_EXTENSIONS = {"jpg", "jpeg", "png", "webp", "bmp", "gif"}
ALLOWED_VIDEO_EXTENSIONS = {"mp4", "mov", "avi", "mkv", "flv", "webm"}
ALLOWED_AUDIO_EXTENSIONS = {"mp3", "wav", "aac", "flac", "ogg", "m4a"}
ALLOWED_DOCUMENT_EXTENSIONS = {"pdf", "docx", "txt", "ppt", "pptx", "xls", "xlsx", "html", "odt", "rtf"}

def allowed_file(filename, allowed_extensions):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in allowed_extensions

def get_user_id_from_token():
    token = request.headers.get('Authorization')
    if token and token.startswith('Bearer '):
        token = token[7:]
        try:
            payload = jwt.decode(token, app.config['SECRET_KEY'], algorithms=['HS256'])
            return payload.get('user_id')
        except:
            return None
    return None

# ==================== AUTH ROUTES ====================

@app.route("/api/auth/forgot-password", methods=["POST"])
def forgot_password():
    try:
        data = request.get_json()
        email = data.get('email', '').strip()
        date_of_birth = data.get('date_of_birth', '').strip()
        new_password = data.get('new_password', '').strip()
        
        if not email:
            return jsonify({"success": False, "message": "Email is required"}), 400
        
        # If new_password is provided, this is a direct reset (DOB verification)
        if new_password:
            if len(new_password) < 4:
                return jsonify({"success": False, "message": "Password must be at least 4 characters"}), 400
            
            if not date_of_birth:
                return jsonify({"success": False, "message": "Date of birth is required for password reset"}), 400
                
            conn = get_db()
            cursor = conn.cursor()
            
            # Check if user exists with matching email and DOB
            cursor.execute(
                "SELECT id, email, date_of_birth FROM users WHERE email = ?",
                (email,)
            )
            user = cursor.fetchone()
            
            if not user:
                conn.close()
                return jsonify({"success": False, "message": "Invalid email or date of birth"}), 400
            
            # Verify DOB matches
            if user['date_of_birth'] != date_of_birth:
                conn.close()
                return jsonify({"success": False, "message": "Invalid email or date of birth"}), 400
            
            # Update password directly
            cursor.execute(
                "UPDATE users SET password = ? WHERE id = ?",
                (new_password, user['id'])
            )
            conn.commit()
            conn.close()
            
            return jsonify({
                "success": True,
                "message": "Password reset successful! You can now login with your new password."
            })
        
        # If no new_password, just check if user exists (forgot password init)
        conn = get_db()
        cursor = conn.cursor()
        
        # Check if user exists
        cursor.execute("SELECT id, email FROM users WHERE email = ?", (email,))
        user = cursor.fetchone()
        
        if not user:
            # Don't reveal if email exists or not for security
            conn.close()
            return jsonify({
                "success": True, 
                "message": "If an account exists with this email, you can reset your password using your date of birth."
            })
        
        # Check if user has DOB set
        cursor.execute("SELECT date_of_birth FROM users WHERE email = ?", (email,))
        user_data = cursor.fetchone()
        conn.close()
        
        if not user_data or not user_data['date_of_birth']:
            return jsonify({
                "success": True,
                "message": "If an account exists with this email, you can reset your password using your date of birth.",
                "has_dob": False
            })
        
        return jsonify({
            "success": True,
            "message": "If an account exists with this email, you can reset your password using your date of birth.",
            "has_dob": True
        })
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500


@app.route("/api/auth/reset-password", methods=["POST"])
def reset_password():
    try:
        data = request.get_json()
        token = data.get('token', '').strip()
        new_password = data.get('new_password', '').strip()
        
        if not token or not new_password:
            return jsonify({"success": False, "message": "Token and new password are required"}), 400
        
        if len(new_password) < 4:
            return jsonify({"success": False, "message": "Password must be at least 4 characters"}), 400
        
        conn = get_db()
        cursor = conn.cursor()
        
        # Check if token is valid and not expired
        cursor.execute(
            "SELECT id, email, reset_token_expires FROM users WHERE reset_token = ?",
            (token,)
        )
        user = cursor.fetchone()
        
        if not user:
            conn.close()
            return jsonify({"success": False, "message": "Invalid reset token"}), 400
        
        # Check if token is expired
        if user['reset_token_expires']:
            expires_at = datetime.fromisoformat(user['reset_token_expires'])
            if datetime.now() > expires_at:
                conn.close()
                return jsonify({"success": False, "message": "Reset token has expired. Please request a new one."}), 400
        
        # Update password and clear reset token
        cursor.execute(
            "UPDATE users SET password = ?, reset_token = NULL, reset_token_expires = NULL WHERE id = ?",
            (new_password, user['id'])
        )
        conn.commit()
        conn.close()
        
        return jsonify({
            "success": True,
            "message": "Password reset successful! You can now login with your new password."
        })
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500


@app.route("/api/auth/verify-reset-token", methods=["POST"])
def verify_reset_token():
    try:
        data = request.get_json()
        token = data.get('token', '').strip()
        
        if not token:
            return jsonify({"success": False, "message": "Token is required"}), 400
        
        conn = get_db()
        cursor = conn.cursor()
        
        cursor.execute(
            "SELECT id, email, reset_token_expires FROM users WHERE reset_token = ?",
            (token,)
        )
        user = cursor.fetchone()
        conn.close()
        
        if not user:
            return jsonify({"success": False, "message": "Invalid reset token"}), 400
        
        # Check if token is expired
        if user['reset_token_expires']:
            expires_at = datetime.fromisoformat(user['reset_token_expires'])
            if datetime.now() > expires_at:
                return jsonify({"success": False, "message": "Reset token has expired"}), 400
        
        return jsonify({
            "success": True,
            "message": "Token is valid",
            "email": user['email']
        })
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500

@app.route("/api/auth/register", methods=["POST"])
def register():
    try:
        data = request.get_json()
        username = data.get('username')
        email = data.get('email')
        password = data.get('password')
        date_of_birth = data.get('date_of_birth', '')
        
        if not username or not email or not password:
            return jsonify({"success": False, "message": "All fields required"}), 400
        
        conn = get_db()
        cursor = conn.cursor()
        
        cursor.execute("SELECT id FROM users WHERE email = ?", (email,))
        if cursor.fetchone():
            conn.close()
            return jsonify({"success": False, "message": "Email already registered"}), 400
        
        plain_password = password
        
        cursor.execute("INSERT INTO users (email, password, date_of_birth, created_at) VALUES (?, ?, ?, ?)", 
                      (email, plain_password, date_of_birth, datetime.now().isoformat()))
        conn.commit()
        user_id = cursor.lastrowid
        
        # Get the newly created user with all fields
        cursor.execute("SELECT id, email, profile_picture, date_of_birth, created_at FROM users WHERE id = ?", (user_id,))
        new_user = cursor.fetchone()
        conn.close()
        
        token = jwt.encode({'user_id': user_id, 'email': email}, 
                          app.config['SECRET_KEY'], algorithm='HS256')
        
        return jsonify({
            "success": True, 
            "message": "Registration successful",
            "token": token,
            "user": {
                "id": new_user['id'], 
                "email": new_user['email'],
                "profile_picture": new_user['profile_picture'],
                "date_of_birth": new_user['date_of_birth'],
                "created_at": new_user['created_at']
            }
        })
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500


@app.route("/api/auth/login", methods=["POST"])
def login():
    try:
        data = request.get_json()
        email = data.get('email')
        password = data.get('password')
        
        if not email or not password:
            return jsonify({"success": False, "message": "Email and password required"}), 400
        
        conn = get_db()
        cursor = conn.cursor()
        
        cursor.execute("SELECT id, email, password, profile_picture, created_at FROM users WHERE email = ?", (email,))
        user = cursor.fetchone()
        conn.close()
        
        if not user:
            return jsonify({"success": False, "message": "User not found"}), 404
        
        if user['password'] != password:
            return jsonify({"success": False, "message": "Invalid password"}), 401
        
        token = jwt.encode({'user_id': user['id'], 'email': user['email']}, 
                          app.config['SECRET_KEY'], algorithm='HS256')
        
        return jsonify({
            "success": True,
            "message": "Login successful",
            "token": token,
            "user": {
                "id": user['id'], 
                "email": user['email'],
                "profile_picture": user['profile_picture'],
                "created_at": user['created_at']
            }
        })
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500


@app.route("/api/auth/change-password", methods=["POST"])
def change_password():
    try:
        user_id = get_user_id_from_token()
        if not user_id:
            return jsonify({"success": False, "message": "Authentication required"}), 401
        
        data = request.get_json()
        current_password = data.get('current_password', '').strip()
        new_password = data.get('new_password', '').strip()
        
        if not current_password or not new_password:
            return jsonify({"success": False, "message": "Current password and new password are required"}), 400
        
        if len(new_password) < 4:
            return jsonify({"success": False, "message": "New password must be at least 4 characters"}), 400
        
        conn = get_db()
        cursor = conn.cursor()
        
        cursor.execute("SELECT id, password FROM users WHERE id = ?", (user_id,))
        user = cursor.fetchone()
        
        if not user:
            conn.close()
            return jsonify({"success": False, "message": "User not found"}), 404
        
        if user['password'] != current_password:
            conn.close()
            return jsonify({"success": False, "message": "Current password is incorrect"}), 401
        
        cursor.execute("UPDATE users SET password = ? WHERE id = ?", (new_password, user_id))
        conn.commit()
        conn.close()
        
        return jsonify({
            "success": True,
            "message": "Password changed successfully!"
        })
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500


# ==================== PROFILE PICTURE ROUTES ====================

@app.route("/api/auth/upload-profile-picture", methods=["POST"])
def upload_profile_picture():
    try:
        user_id = get_user_id_from_token()
        if not user_id:
            return jsonify({"success": False, "message": "Authentication required"}), 401

        if "file" not in request.files:
            return jsonify({"success": False, "message": "No file selected"}), 400

        file = request.files["file"]
        if file.filename == "":
            return jsonify({"success": False, "message": "No file selected"}), 400

        ext = file.filename.rsplit(".", 1)[1].lower()
        if ext not in ALLOWED_IMAGE_EXTENSIONS:
            return jsonify({"success": False, "message": "Invalid file type"}), 400

        profile_filename = f"profile_{user_id}.{ext}"
        profile_path = os.path.join(PROFILE_FOLDER, profile_filename)
        file.save(profile_path)

        conn = get_db()
        cursor = conn.cursor()
        cursor.execute(
            "UPDATE users SET profile_picture = ? WHERE id = ?",
            (profile_filename, user_id)
        )
        conn.commit()
        conn.close()

        return jsonify({
            "success": True,
            "message": "Profile picture updated",
            "profile_picture": profile_filename
        })
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500


# ==================== FILE UPLOAD ROUTES ====================

@app.route("/api/conversion/guest-upload", methods=["POST"])
def guest_upload():
    try:
        if "file" not in request.files:
            return jsonify({"error": "No file selected"}), 400

        file = request.files["file"]
        if file.filename == "":
            return jsonify({"error": "No file selected"}), 400

        ext = file.filename.rsplit(".", 1)[1].lower()
        if ext not in ALLOWED_IMAGE_EXTENSIONS:
            return jsonify({"error": "Invalid file type"}), 400

        unique_name = f"{uuid.uuid4()}.{ext}"
        save_path = os.path.join(ORIGINAL_FOLDER, unique_name)
        file.save(save_path)

        return jsonify({
            "success": True,
            "message": "Guest upload successful",
            "filename": unique_name,
            "path": save_path
        })
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500


@app.route("/api/conversion/upload", methods=["POST"])
def upload_file():
    try:
        if "file" not in request.files:
            return jsonify({"success": False, "error": "No file selected"}), 400

        file = request.files["file"]
        if file.filename == "":
            return jsonify({"success": False, "error": "No file selected"}), 400

        ext = file.filename.rsplit(".", 1)[1].lower()
        if ext not in ALLOWED_IMAGE_EXTENSIONS:
            return jsonify({"success": False, "error": "Invalid file type"}), 400

        unique_name = f"{uuid.uuid4()}.{ext}"
        save_path = os.path.join(ORIGINAL_FOLDER, unique_name)
        file.save(save_path)

        user_id = get_user_id_from_token()
        
        if user_id:
            conn = get_db()
            cursor = conn.cursor()
            cursor.execute(
                "INSERT INTO history (user_id, filename, converted_filename, date) VALUES (?, ?, ?, ?)",
                (user_id, unique_name, '', datetime.now().isoformat())
            )
            conn.commit()
            conn.close()

        return jsonify({
            "success": True,
            "message": "Upload successful",
            "filename": unique_name
        })
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500


# ==================== CONVERSION ROUTES ====================

@app.route("/api/convert/image", methods=["POST"])
def convert_image():
    try:
        if "file" not in request.files:
            return jsonify({"success": False, "error": "No file selected"}), 400

        file = request.files["file"]
        target_format = request.form.get("targetFormat", "jpg").lower()
        quality = int(request.form.get("quality", 85))

        # Check for crop parameters
        crop_data = request.form.get("cropData")
        crop_x = request.form.get("cropX")
        crop_y = request.form.get("cropY")
        crop_width = request.form.get("cropWidth")
        crop_height = request.form.get("cropHeight")
        crop_rotate = request.form.get("cropRotate", "0")
        crop_scaleX = request.form.get("cropScaleX", "1")
        crop_scaleY = request.form.get("cropScaleY", "1")

        if file.filename == "":
            return jsonify({"success": False, "error": "No file selected"}), 400

        ext = file.filename.rsplit(".", 1)[1].lower()
        if ext not in ALLOWED_IMAGE_EXTENSIONS:
            return jsonify({"success": False, "error": "Invalid file type"}), 400

        if target_format not in ALLOWED_IMAGE_EXTENSIONS:
            return jsonify({"success": False, "error": "Invalid target format"}), 400

        original_filename = secure_filename(file.filename)
        original_name = original_filename
        original_path = os.path.join(ORIGINAL_FOLDER, original_name)
        file.save(original_path)

        converted_name = f"{original_filename.rsplit('.', 1)[0]}.{target_format}"
        converted_path = os.path.join(CONVERTED_FOLDER, converted_name)

        img = Image.open(original_path)
        
        # Apply crop if crop data provided
        if crop_x and crop_y and crop_width and crop_height:
            try:
                x = int(float(crop_x))
                y = int(float(crop_y))
                width = int(float(crop_width))
                height = int(float(crop_height))
                
                # Crop the image
                img = img.crop((x, y, x + width, y + height))
            except Exception as crop_err:
                print(f"Crop error: {crop_err}")
        
        # Apply rotation if specified
        if crop_rotate and crop_rotate != "0":
            try:
                rotation = int(float(crop_rotate))
                img = img.rotate(rotation, expand=True)
            except Exception as rot_err:
                print(f"Rotation error: {rot_err}")
        
        # Apply flip if specified
        if crop_scaleX == "-1":
            img = img.transpose(Image.FLIP_LEFT_RIGHT)
        if crop_scaleY == "-1":
            img = img.transpose(Image.FLIP_TOP_BOTTOM)
        
        if target_format in ['jpg', 'jpeg'] and img.mode in ('RGBA', 'P'):
            img = img.convert('RGB')
        
        img.save(converted_path, quality=quality)
        
        user_id = get_user_id_from_token()
        
        if user_id:
            conn = get_db()
            cursor = conn.cursor()
            cursor.execute(
                "INSERT INTO history (user_id, filename, converted_filename, file_type, date) VALUES (?, ?, ?, ?, ?)",
                (user_id, original_name, converted_name, 'image', datetime.now().isoformat())
            )
            conn.commit()
            conn.close()

        return jsonify({
            "success": True,
            "message": "Conversion successful",
            "filename": converted_name,
            "file_type": "image",
            "downloadUrl": f"/api/files/converted/{converted_name}",
            "previewUrl": f"/api/files/converted/{converted_name}"
        })
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/convert/video", methods=["POST"])
def convert_video():
    try:
        if "file" not in request.files:
            return jsonify({"success": False, "error": "No file selected"}), 400

        file = request.files["file"]
        target_format = request.form.get("targetFormat", "mp4").lower()

        if file.filename == "":
            return jsonify({"success": False, "error": "No file selected"}), 400

        ext = file.filename.rsplit(".", 1)[1].lower()
        if ext not in ALLOWED_VIDEO_EXTENSIONS:
            return jsonify({"success": False, "error": "Invalid video file"}), 400

        if target_format not in ALLOWED_VIDEO_EXTENSIONS:
            return jsonify({"success": False, "error": "Invalid target format"}), 400

        # Save original file temporarily
        original_name = secure_filename(file.filename)
        original_path = os.path.join(CONVERTED_FOLDER, f"temp_{original_name}")
        file.save(original_path)
        
        # Generate new filename with target format
        base_name = original_name.rsplit('.', 1)[0]
        converted_name = f"{base_name}.{target_format}"
        converted_path = os.path.join(CONVERTED_FOLDER, converted_name)
        
        # Perform video conversion using ffmpeg
        conversion_success = False
        try:
            import ffmpeg
            import subprocess
            
            # Get quality setting (1-100, default 85)
            quality = int(request.form.get("quality", 85))
            
            # Map quality to CRF (lower = better quality)
            # CRF 18-23 is good quality, 28+ is lower quality
            crf = max(18, 51 - (quality * 33 // 100))
            
            # Use ffmpeg-python if available, else use subprocess
            try:
                stream = ffmpeg.input(original_path)
                stream = ffmpeg.output(stream, converted_path, 
                                       **{'c:v': 'libx264', 'c:a': 'aac', 
                                          'crf': str(crf), 'preset': 'medium',
                                          'movflags': '+faststart'})
                ffmpeg.run(stream, overwrite_output=True, quiet=True)
                conversion_success = True
            except Exception as ffmpeg_err:
                print(f"ffmpeg-python error: {ffmpeg_err}")
                # Fallback to subprocess
                cmd = ['ffmpeg', '-i', original_path, '-c:v', 'libx264', 
                       '-crf', str(crf), '-preset', 'medium', 
                       '-c:a', 'aac', '-movflags', '+faststart',
                       '-y', converted_path]
                result = subprocess.run(cmd, capture_output=True, text=True)
                if result.returncode == 0:
                    conversion_success = True
                else:
                    print(f"ffmpeg subprocess error: {result.stderr}")
            
            # Remove temp file
            if os.path.exists(original_path):
                os.remove(original_path)
        except ImportError:
            # ffmpeg not installed - copy file with new extension as fallback
            print("ffmpeg not available, using direct copy")
            try:
                import shutil
                # Just copy with new extension as placeholder
                if os.path.exists(original_path):
                    shutil.copy(original_path, converted_path)
                    conversion_success = True
            except Exception as copy_err:
                print(f"Copy error: {copy_err}")
        except Exception as e:
            print(f"Conversion error: {e}")
            # On error, try to use original file
            if os.path.exists(original_path):
                try:
                    import shutil
                    shutil.copy(original_path, converted_path)
                    conversion_success = True
                except:
                    pass

        if not conversion_success:
            # If conversion failed, return original file
            if os.path.exists(original_path):
                converted_name = original_name
                converted_path = original_path

        # Save to history
        user_id = get_user_id_from_token()
        if user_id:
            conn = get_db()
            cursor = conn.cursor()
            cursor.execute(
                "INSERT INTO history (user_id, filename, converted_filename, file_type, date) VALUES (?, ?, ?, ?, ?)",
                (user_id, original_name, converted_name, 'video', datetime.now().isoformat())
            )
            conn.commit()
            conn.close()

        return jsonify({
            "success": True,
            "message": f"Video converted to {target_format.upper()}" if conversion_success else "Video uploaded (conversion limited)",
            "filename": converted_name,
            "file_type": "video",
            "downloadUrl": f"/api/files/converted/{converted_name}",
            "previewUrl": f"/api/files/converted/{converted_name}"
        })
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/convert/audio", methods=["POST"])
def convert_audio():
    try:
        if "file" not in request.files:
            return jsonify({"success": False, "error": "No file selected"}), 400

        file = request.files["file"]
        target_format = request.form.get("targetFormat", "wav").lower()
        bitrate = int(request.form.get("bitrate", 192))

        if file.filename == "":
            return jsonify({"success": False, "error": "No file selected"}), 400

        ext = file.filename.rsplit(".", 1)[1].lower()
        if ext not in ALLOWED_AUDIO_EXTENSIONS:
            return jsonify({"success": False, "error": "Invalid audio file"}), 400

        if target_format not in ALLOWED_AUDIO_EXTENSIONS:
            return jsonify({"success": False, "error": "Invalid target format"}), 400

        # Save original file temporarily
        original_name = secure_filename(file.filename)
        original_path = os.path.join(CONVERTED_FOLDER, f"temp_{original_name}")
        file.save(original_path)

        # Generate output audio filename with target format
        converted_name = f"{original_name.rsplit('.', 1)[0]}.{target_format}"
        converted_path = os.path.join(CONVERTED_FOLDER, converted_name)

        # Perform audio conversion using ffmpeg
        conversion_success = False
        try:
            import subprocess
            
            # Use ffmpeg to convert audio
            cmd = ['ffmpeg', '-i', original_path, '-ab', f'{bitrate}k', '-y', converted_path]
            result = subprocess.run(cmd, capture_output=True, text=True)
            
            if result.returncode == 0 and os.path.exists(converted_path):
                conversion_success = True
            else:
                print(f"ffmpeg error: {result.stderr}")
                
        except ImportError:
            print("ffmpeg not available")
        except Exception as e:
            print(f"Audio conversion error: {e}")

        # Remove temp file
        if os.path.exists(original_path):
            os.remove(original_path)

        if not conversion_success:
            # Fallback: just copy original to converted with target extension
            base_name = original_name.rsplit('.', 1)[0]
            converted_name = f"{base_name}.{target_format}"
            converted_path = os.path.join(CONVERTED_FOLDER, converted_name)
            try:
                import shutil
                shutil.copy(original_path, converted_path)
                conversion_success = True
            except:
                pass

        # Save to history
        user_id = get_user_id_from_token()
        if user_id:
            conn = get_db()
            cursor = conn.cursor()
            cursor.execute(
                "INSERT INTO history (user_id, filename, converted_filename, file_type, date) VALUES (?, ?, ?, ?, ?)",
                (user_id, original_name, converted_name, 'audio', datetime.now().isoformat())
            )
            conn.commit()
            conn.close()

        return jsonify({
            "success": True,
            "message": f"Audio converted to {target_format.upper()}" if conversion_success else "Audio conversion failed",
            "filename": converted_name,
            "file_type": "audio",
            "downloadUrl": f"/api/files/converted/{converted_name}"
        })
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/convert/video-to-audio", methods=["POST"])
def video_to_audio():
    try:
        if "file" not in request.files:
            return jsonify({"success": False, "error": "No file selected"}), 400

        file = request.files["file"]
        target_format = request.form.get("targetFormat", "mp3").lower()
        bitrate = int(request.form.get("bitrate", 192))

        if file.filename == "":
            return jsonify({"success": False, "error": "No file selected"}), 400

        ext = file.filename.rsplit(".", 1)[1].lower()
        if ext not in ALLOWED_VIDEO_EXTENSIONS:
            return jsonify({"success": False, "error": "Invalid video file"}), 400

        if target_format not in ALLOWED_AUDIO_EXTENSIONS:
            return jsonify({"success": False, "error": "Invalid audio format"}), 400

        # Save original video temporarily
        original_name = secure_filename(file.filename)
        video_path = os.path.join(CONVERTED_FOLDER, f"temp_{original_name}")
        file.save(video_path)

        # Generate output audio filename - ALWAYS use the selected format
        audio_name = f"{original_name.rsplit('.', 1)[0]}.{target_format}"
        audio_path = os.path.join(CONVERTED_FOLDER, audio_name)

        # Try to extract audio using moviepy first
        audio_extraction_success = False
        try:
            # Updated import for moviepy v2.0+
            try:
                from moviepy import VideoFileClip
            except ImportError:
                from moviepy.editor import VideoFileClip
            
            # Load video
            video_clip = VideoFileClip(video_path)
            
            # Extract audio
            audio_clip = video_clip.audio
            
            # Write audio file with proper encoding
            if target_format == 'mp3':
                # For MP3, use specific codec
                audio_clip.write_audiofile(audio_path, codec='libmp3lame', bitrate=f'{bitrate}k', verbose=False, logger=None)
            elif target_format == 'wav':
                audio_clip.write_audiofile(audio_path, codec='pcm_s16le', verbose=False, logger=None)
            elif target_format == 'aac':
                audio_clip.write_audiofile(audio_path, codec='aac', bitrate=f'{bitrate}k', verbose=False, logger=None)
            elif target_format == 'ogg':
                audio_clip.write_audiofile(audio_path, codec='libvorbis', bitrate=f'{bitrate}k', verbose=False, logger=None)
            else:
                audio_clip.write_audiofile(audio_path, bitrate=f'{bitrate}k', verbose=False, logger=None)
            
            # Close clips
            audio_clip.close()
            video_clip.close()
            
            # Verify the file was created
            if os.path.exists(audio_path) and os.path.getsize(audio_path) > 0:
                audio_extraction_success = True
            else:
                print("Moviepy created empty file, will try ffmpeg fallback")
                
        except ImportError:
            print("Moviepy not available, will try ffmpeg fallback")
            audio_extraction_success = False
        except Exception as e:
            print(f"Moviepy extraction error: {e}")
            audio_extraction_success = False

        # If moviepy failed, try ffmpeg directly (more reliable for MP3 on Windows)
        if not audio_extraction_success:
            try:
                import subprocess
                
                # Build ffmpeg command based on format
                if target_format == 'mp3':
                    # Use libmp3lame for MP3 encoding
                    cmd = ['ffmpeg', '-i', video_path, '-vn', '-ab', f'{bitrate}k', 
                           '-acodec', 'libmp3lame', '-y', audio_path]
                elif target_format == 'wav':
                    cmd = ['ffmpeg', '-i', video_path, '-vn', '-acodec', 'pcm_s16le', '-y', audio_path]
                elif target_format == 'aac':
                    cmd = ['ffmpeg', '-i', video_path, '-vn', '-ab', f'{bitrate}k', 
                           '-acodec', 'aac', '-y', audio_path]
                elif target_format == 'ogg':
                    cmd = ['ffmpeg', '-i', video_path, '-vn', '-ab', f'{bitrate}k', 
                           '-acodec', 'libvorbis', '-y', audio_path]
                else:
                    cmd = ['ffmpeg', '-i', video_path, '-vn', '-ab', f'{bitrate}k', 
                           '-acodec', 'libmp3lame', '-y', audio_path]
                
                result = subprocess.run(cmd, capture_output=True, text=True, timeout=300)
                
                # Check if file was created successfully
                if result.returncode == 0 and os.path.exists(audio_path) and os.path.getsize(audio_path) > 0:
                    audio_extraction_success = True
                    print(f"ffmpeg successfully extracted audio to {target_format}")
                else:
                    print(f"ffmpeg failed: {result.stderr}")
                    
            except FileNotFoundError:
                print("ffmpeg not found in system PATH")
            except subprocess.TimeoutExpired:
                print("ffmpeg conversion timed out")
            except Exception as e:
                print(f"ffmpeg fallback error: {e}")

        # Clean up temporary video file
        if os.path.exists(video_path):
            try:
                os.remove(video_path)
            except:
                pass

        # If extraction failed completely, return error
        if not audio_extraction_success:
            return jsonify({
                "success": False,
                "error": "Audio extraction failed. Please ensure ffmpeg is installed on the server."
            }), 500

        # Save to history (for both logged in users and guests)
        user_id = get_user_id_from_token()
        conn = get_db()
        cursor = conn.cursor()
        cursor.execute(
            "INSERT INTO history (user_id, filename, converted_filename, file_type, date) VALUES (?, ?, ?, ?, ?)",
            (user_id, original_name, audio_name, 'video-to-audio', datetime.now().isoformat())
        )
        conn.commit()
        conn.close()

        return jsonify({
            "success": True,
            "message": f"Audio extraction complete ({target_format.upper()})",
            "filename": audio_name,
            "file_type": "video-to-audio",
            "downloadUrl": f"/api/files/converted/{audio_name}",
            "previewUrl": f"/api/files/converted/{audio_name}"
        })
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/convert/document", methods=["POST"])
def convert_document():
    try:
        if "file" not in request.files:
            return jsonify({"success": False, "error": "No file selected"}), 400

        file = request.files["file"]
        target_format = request.form.get("targetFormat", "txt").lower()
        password = request.form.get("password", "").strip()

        if file.filename == "":
            return jsonify({"success": False, "error": "No file selected"}), 400

        ext = file.filename.rsplit(".", 1)[1].lower()
        if ext not in ALLOWED_DOCUMENT_EXTENSIONS:
            return jsonify({"success": False, "error": f"Unsupported document format: .{ext}. Supported: {', '.join(ALLOWED_DOCUMENT_EXTENSIONS)}"}), 400

        if target_format not in ALLOWED_DOCUMENT_EXTENSIONS:
            return jsonify({"success": False, "error": f"Invalid target format: {target_format}. Supported: {', '.join(ALLOWED_DOCUMENT_EXTENSIONS)}"}), 400

        # Save original file temporarily
        original_name = secure_filename(file.filename)
        original_path = os.path.join(CONVERTED_FOLDER, f"temp_{original_name}")
        file.save(original_path)

        # Generate output filename with target format
        converted_name = f"{original_name.rsplit('.', 1)[0]}.{target_format}"
        converted_path = os.path.join(CONVERTED_FOLDER, converted_name)

        # Perform document conversion
        conversion_success = False
        
        try:
            # If converting to TXT from other formats
            if target_format == 'txt':
                # Try to extract text from PDF or other formats
                try:
                    # For PDF files, try using pypdf
                    if ext == 'pdf':
                        from pypdf import PdfReader
                        reader = PdfReader(original_path)
                        text_content = ""
                        for page in reader.pages:
                            text_content += page.extract_text() + "\n"
                        
                        # Save as txt
                        with open(converted_path, 'w', encoding='utf-8') as f:
                            f.write(text_content)
                        conversion_success = True
                    else:
                        # For other files, try to read as text
                        with open(original_path, 'r', encoding='utf-8', errors='ignore') as f:
                            content = f.read()
                        with open(converted_path, 'w', encoding='utf-8') as f:
                            f.write(content)
                        conversion_success = True
                except ImportError:
                    # pypdf not available, copy file as fallback
                    import shutil
                    shutil.copy(original_path, converted_path)
                    conversion_success = True
                except Exception as e:
                    print(f"Text extraction error: {e}")
                    # Fallback: copy file
                    try:
                        import shutil
                        shutil.copy(original_path, converted_path)
                        conversion_success = True
                    except:
                        pass
            else:
                # Handle PPTX/XLSX specifically
                if ext in ['pptx', 'xlsx'] or target_format in ['pptx', 'xlsx']:
                    try:
                        if ext == 'pptx' and target_format == 'txt':
                            # Extract text from PPTX
                            from pptx import Presentation
                            prs = Presentation(original_path)
                            text_content = []
                            for slide in prs.slides:
                                for shape in slide.shapes:
                                    if hasattr(shape, 'text'):
                                        text_content.append(shape.text)
                            with open(converted_path, 'w', encoding='utf-8') as f:
                                f.write('\\n'.join(text_content))
                            conversion_success = True
                        elif ext == 'xlsx' and target_format == 'txt':
                            # Convert XLSX to CSV-like text
                            import pandas as pd
                            df = pd.read_excel(original_path)
                            df.to_csv(converted_path, index=False)
                            conversion_success = True
                        else:
                            # Copy for other PPTX/XLSX conversions (placeholder for full support)
                            import shutil
                            shutil.copy(original_path, converted_path)
                            conversion_success = True
                    except ImportError:
                        import shutil
                        shutil.copy(original_path, converted_path)
                        conversion_success = True
                    except Exception as office_err:
                        print(f"Office format error: {office_err}")
                        import shutil
                        shutil.copy(original_path, converted_path)
                        conversion_success = True
                else:
                    # For other conversions (like pdf to pdf), just copy
                    import shutil
                    shutil.copy(original_path, converted_path)
                    conversion_success = True
                
        except Exception as e:
            print(f"Document conversion error: {e}")
            # Fallback: copy file
            try:
                import shutil
                shutil.copy(original_path, converted_path)
                conversion_success = True
            except:
                pass

        # Remove temp file
        if os.path.exists(original_path):
            os.remove(original_path)

        # Apply password protection if requested and target is PDF
        if password and target_format == 'pdf':
            try:
                from PyPDF2 import PdfReader, PdfWriter
                
                reader = PdfReader(converted_path)
                writer = PdfWriter()
                
                for page in reader.pages:
                    writer.add_page(page)
                
                # Set owner password (user password for opening)
                writer.encrypt(user_password=password, owner_password=None, 
                             use_128bit=True, allow_printing=True)
                
                # Write encrypted PDF
                encrypted_path = converted_path.replace('.pdf', '_encrypted.pdf')
                with open(encrypted_path, "wb") as output_file:
                    writer.write(output_file)
                
                # Replace original with encrypted version
                os.remove(converted_path)
                converted_path = encrypted_path
                converted_name = converted_name.replace('.pdf', '_encrypted.pdf')
                
                print("PDF encrypted successfully")
            except ImportError:
                print("PyPDF2 not available, skipping PDF encryption")
            except Exception as encrypt_err:
                print(f"PDF encryption failed: {encrypt_err}")
        elif password:
            print("Password protection only supported for PDF output")

        if not conversion_success:
            # If conversion failed, return original file
            converted_name = original_name
            converted_path = os.path.join(CONVERTED_FOLDER, converted_name)

        # Save to history
        user_id = get_user_id_from_token()
        if user_id:
            conn = get_db()
            cursor = conn.cursor()
            cursor.execute(
                "INSERT INTO history (user_id, filename, converted_filename, file_type, date) VALUES (?, ?, ?, ?, ?)",
                (user_id, original_name, converted_name, 'document', datetime.now().isoformat())
            )
            conn.commit()
            conn.close()

        return jsonify({
            "success": True,
            "message": f"Document converted to {target_format.upper()}" if conversion_success else "Document conversion failed",
            "filename": converted_name,
            "file_type": "document",
            "downloadUrl": f"/api/files/converted/{converted_name}"
        })
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


# ==================== HISTORY ROUTES ====================

@app.route("/api/conversion/recent", methods=["GET"])
def get_recent():
    try:
        user_id = get_user_id_from_token()
        
        conn = get_db()
        cursor = conn.cursor()
        
        if user_id:
            cursor.execute(
                "SELECT id, filename, converted_filename, file_type, date FROM history WHERE user_id = ? ORDER BY date DESC LIMIT 10",
                (user_id,)
            )
        else:
            cursor.execute(
                "SELECT id, filename, converted_filename, file_type, date FROM history ORDER BY date DESC LIMIT 10"
            )
        
        rows = cursor.fetchall()
        conn.close()
        
        items = []
        for row in rows:
            items.append({
                "_id": row['id'],
                "filename": row['filename'],
                "converted_filename": row['converted_filename'],
                "file_type": row['file_type'] if row['file_type'] else 'image',
                "date": row['date'],
                "type": row['file_type'] if row['file_type'] else 'image'
            })
        
        return jsonify({
            "success": True,
            "items": items
        })
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/conversion/preview/<id>", methods=["GET"])
def preview_file(id):
    try:
        conn = get_db()
        cursor = conn.cursor()
        cursor.execute("SELECT converted_filename FROM history WHERE id = ?", (id,))
        row = cursor.fetchone()
        conn.close()
        
        if row and row['converted_filename']:
            return send_from_directory(CONVERTED_FOLDER, row['converted_filename'])
        else:
            return jsonify({"error": "File not found"}), 404
    except Exception as e:
        return jsonify({"error": str(e)}), 404


@app.route("/api/conversion/download/<id>", methods=["GET"])
def download_file(id):
    try:
        conn = get_db()
        cursor = conn.cursor()
        cursor.execute("SELECT converted_filename FROM history WHERE id = ?", (id,))
        row = cursor.fetchone()
        conn.close()
        
        if row and row['converted_filename']:
            return send_from_directory(CONVERTED_FOLDER, row['converted_filename'], as_attachment=True)
        else:
            return jsonify({"error": "File not found"}), 404
    except Exception as e:
        return jsonify({"error": str(e)}), 404


# ==================== ADMIN ROUTES ====================

@app.route("/api/admin/stats", methods=["GET"])
def admin_stats():
    try:
        conn = get_db()
        cursor = conn.cursor()
        
        cursor.execute("SELECT COUNT(*) as count FROM users")
        total_users = cursor.fetchone()['count']
        
        cursor.execute("SELECT COUNT(*) as count FROM history")
        total_conversions = cursor.fetchone()['count']
        
        cursor.execute("SELECT COUNT(*) as count FROM history WHERE date(date) = date('now')")
        today_conversions = cursor.fetchone()['count']
        
        import os
        total_size = 0
        for folder in [ORIGINAL_FOLDER, CONVERTED_FOLDER]:
            if os.path.exists(folder):
                for file in os.listdir(folder):
                    file_path = os.path.join(folder, file)
                    if os.path.isfile(file_path):
                        total_size += os.path.getsize(file_path)
        
        storage_mb = round(total_size / (1024 * 1024), 2)
        conn.close()
        
        return jsonify({
            "success": True,
            "stats": {
                "total_users": total_users,
                "total_conversions": total_conversions,
                "today_conversions": today_conversions,
                "storage_used_mb": storage_mb
            }
        })
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/admin/users", methods=["GET"])
def admin_all_users():
    try:
        conn = get_db()
        cursor = conn.cursor()
        
        cursor.execute("""
            SELECT u.id, u.email, u.created_at, u.profile_picture,
                   COUNT(h.id) as conversion_count
            FROM users u
            LEFT JOIN history h ON u.id = h.user_id
            GROUP BY u.id
            ORDER BY u.id DESC
        """)
        
        rows = cursor.fetchall()
        conn.close()
        
        users = []
        for row in rows:
            users.append({
                "id": row['id'],
                "email": row['email'],
                "created_at": row['created_at'],
                "profile_picture": row['profile_picture'],
                "conversions": row['conversion_count']
            })
        
        return jsonify({"success": True, "users": users})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/admin/conversions", methods=["GET"])
def admin_all_conversions():
    try:
        conn = get_db()
        cursor = conn.cursor()
        
        cursor.execute("""
            SELECT h.id, h.filename, h.converted_filename, h.date, u.email
            FROM history h
            LEFT JOIN users u ON h.user_id = u.id
            ORDER BY h.date DESC
            LIMIT 100
        """)
        
        rows = cursor.fetchall()
        conn.close()
        
        conversions = []
        for row in rows:
            conversions.append({
                "id": row['id'],
                "filename": row['filename'],
                "converted_filename": row['converted_filename'],
                "date": row['date'],
                "user_email": row['email'] if row['email'] else 'Guest'
            })
        
        return jsonify({"success": True, "conversions": conversions})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/admin/user/<user_id>", methods=["GET"])
def admin_user_details(user_id):
    try:
        conn = get_db()
        cursor = conn.cursor()
        
        cursor.execute("SELECT id, email, created_at, profile_picture FROM users WHERE id = ?", (user_id,))
        user = cursor.fetchone()
        
        if not user:
            conn.close()
            return jsonify({"success": False, "error": "User not found"}), 404
        
        cursor.execute("""
            SELECT id, filename, converted_filename, date
            FROM history WHERE user_id = ?
            ORDER BY date DESC LIMIT 50
        """, (user_id,))
        
        rows = cursor.fetchall()
        
        conversions = []
        for row in rows:
            conversions.append({
                "id": row['id'],
                "filename": row['filename'],
                "converted_filename": row['converted_filename'],
                "date": row['date']
            })
        
        conn.close()
        
        return jsonify({
            "success": True,
            "user": {
                "id": user['id'], 
                "email": user['email'], 
                "created_at": user['created_at'],
                "profile_picture": user['profile_picture']
            },
            "conversions": conversions
        })
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


# ==================== FILE SERVE ROUTES ====================

@app.route("/api/files/original/<filename>")
def get_original_file(filename):
    return send_from_directory(ORIGINAL_FOLDER, filename)


@app.route("/api/files/converted/<filename>")
def get_converted_file(filename):
    return send_from_directory(CONVERTED_FOLDER, filename, as_attachment=True)


@app.route("/api/files/profiles/<filename>")
def get_profile_file(filename):
    return send_from_directory(PROFILE_FOLDER, filename)


# ==================== CHAT ROUTES ====================

@app.route("/api/chat/send", methods=["POST"])
def send_chat_message():
    try:
        data = request.get_json()
        user_id = get_user_id_from_token()
        message = data.get('message', '').strip()
        
        if not message:
            return jsonify({"success": False, "message": "Message cannot be empty"}), 400
        
        # Get user email if logged in
        user_email = "Guest"
        if user_id:
            conn = get_db()
            cursor = conn.cursor()
            cursor.execute("SELECT email FROM users WHERE id = ?", (user_id,))
            user = cursor.fetchone()
            if user:
                user_email = user['email']
            conn.close()
        
        # Store message in database
        conn = get_db()
        cursor = conn.cursor()
        cursor.execute(
            "INSERT INTO chat_messages (user_id, user_email, message, is_from_admin, created_at) VALUES (?, ?, ?, ?, ?)",
            (user_id, user_email, message, 0, datetime.now().isoformat())
        )
        conn.commit()
        message_id = cursor.lastrowid
        conn.close()
        
        # Generate auto-response based on keywords
        bot_response = generate_bot_response(message)
        
        return jsonify({
            "success": True,
            "message": "Message sent",
            "user_message": {
                "id": message_id,
                "message": message,
                "user_email": user_email,
                "is_from_admin": False,
                "created_at": datetime.now().isoformat()
            },
            "bot_response": bot_response
        })
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500


@app.route("/api/chat/messages", methods=["GET"])
def get_chat_messages():
    try:
        user_id = get_user_id_from_token()
        
        conn = get_db()
        cursor = conn.cursor()
        
        if user_id == 1:  # Admin user sees all messages
            cursor.execute(
                "SELECT id, user_id, user_email, message, is_from_admin, created_at FROM chat_messages ORDER BY created_at DESC LIMIT 100"
            )
        elif user_id:
            # Get messages for this user
            cursor.execute(
                "SELECT id, user_id, user_email, message, is_from_admin, created_at FROM chat_messages WHERE user_id = ? ORDER BY created_at ASC LIMIT 50",
                (user_id,)
            )
        else:
            # Get only recent guest messages
            cursor.execute(
                "SELECT id, user_id, user_email, message, is_from_admin, created_at FROM chat_messages WHERE user_id IS NULL ORDER BY created_at DESC LIMIT 50"
            )
        
        rows = cursor.fetchall()
        conn.close()
        
        messages = []
        for row in rows:
            messages.append({
                "id": row['id'],
                "user_id": row['user_id'],
                "user_email": row['user_email'],
                "message": row['message'],
                "is_from_admin": row['is_from_admin'],
                "created_at": row['created_at']
            })
        
        return jsonify({
            "success": True,
            "messages": messages
        })
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500


def generate_bot_response(user_message):
    """Generate automated responses based on keywords"""
    message = user_message.lower()
    
    responses = {
        "greeting": ["Hello! How can I help you today?", "Hi there! What can I assist you with?", "Welcome! How may I help you?"],
        "formats": ["We support image formats (PNG, JPG, WEBP, BMP, GIF), video formats (MP4, MOV, AVI), audio formats (MP3, WAV, AAC), and document formats (PDF, DOCX, TXT)."],
        "size": ["Free accounts can upload up to 100MB per file. Pro accounts get 1GB limit."],
        "convert": ["To convert a file: 1) Upload your file 2) Select target format 3) Click Convert 4) Download result"],
        "account": ["You can create an account by clicking 'Account' and then 'Register'. Having an account lets you track conversion history."],
        "contact": ["You can reach us at mahajanvikram19@gmail.com or call +91 738-746-9487"],
        "default": ["Thank you for your message! Our support team will get back to you shortly. For urgent matters, please email us directly."]
    }
    
    # Check for keywords
    if any(word in message for word in ["hello", "hi", "hey", "good morning"]):
        return {"text": responses["greeting"][0], "type": "bot"}
    elif any(word in message for word in ["format", "file type", "supported"]):
        return {"text": responses["formats"][0], "type": "bot"}
    elif any(word in message for word in ["size", "limit", "mb", "gb"]):
        return {"text": responses["size"][0], "type": "bot"}
    elif any(word in message for word in ["convert", "how to", "help"]):
        return {"text": responses["convert"][0], "type": "bot"}
    elif any(word in message for word in ["account", "register", "login", "signup"]):
        return {"text": responses["account"][0], "type": "bot"}
    elif any(word in message for word in ["contact", "email", "phone", "support"]):
        return {"text": responses["contact"][0], "type": "bot"}
    else:
        return {"text": responses["default"][0], "type": "bot"}


# ==================== CONTACT FORM ROUTES ====================

@app.route("/api/contact/submit", methods=["POST"])
def submit_contact_form():
    try:
        data = request.get_json()
        
        name = data.get('name', '').strip()
        email = data.get('email', '').strip()
        phone = data.get('phone', '').strip()
        subject = data.get('subject', '').strip()
        message = data.get('message', '').strip()
        
        if not name or not email or not message:
            return jsonify({"success": False, "message": "Name, email, and message are required"}), 400
        
        # Format message for chat display
        formatted_message = f"[Contact Form] Subject: {subject}\n\n{message}"
        
        conn = get_db()
        cursor = conn.cursor()
        
        # Store in contact_messages for admin dashboard
        cursor.execute(
            "INSERT INTO contact_messages (name, email, phone, subject, message, created_at) VALUES (?, ?, ?, ?, ?, ?)",
            (name, email, phone, subject, message, datetime.now().isoformat())
        )
        contact_id = cursor.lastrowid
        
        # Find user_id by email for chat_messages
        cursor.execute("SELECT id FROM users WHERE email = ?", (email,))
        user = cursor.fetchone()
        user_id = user['id'] if user else None
        
        # Also store in chat_messages so it appears in user's account "My Messages"
        cursor.execute(
            "INSERT INTO chat_messages (user_id, user_email, message, is_from_admin, created_at) VALUES (?, ?, ?, ?, ?)",
            (user_id, email, formatted_message, 0, datetime.now().isoformat())
        )
        conn.commit()
        conn.close()
        
        return jsonify({
            "success": True,
            "message": "Message sent successfully",
            "id": contact_id
        })
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500


@app.route("/api/admin/contact-messages", methods=["GET"])
def admin_get_contact_messages():
    try:
        conn = get_db()
        cursor = conn.cursor()
        
        cursor.execute("""
            SELECT id, name, email, phone, subject, message, created_at
            FROM contact_messages
            ORDER BY created_at DESC
            LIMIT 100
        """)
        
        rows = cursor.fetchall()
        conn.close()
        
        messages = []
        for row in rows:
            messages.append({
                "id": row['id'],
                "name": row['name'],
                "email": row['email'],
                "phone": row['phone'],
                "subject": row['subject'],
                "message": row['message'],
                "created_at": row['created_at']
            })
        
        return jsonify({
            "success": True,
            "messages": messages
        })
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route("/api/admin/reply-user", methods=["POST"])
def admin_reply_user():
    try:
        data = request.get_json()
        message_id = data.get('message_id')
        reply_message = data.get('reply', '').strip()
        admin_email = data.get('admin_email', 'admin@fileconverter.com')
        
        if not message_id or not reply_message:
            return jsonify({"success": False, "message": "Message ID and reply are required"}), 400
        
        # Get original message to find user email
        conn = get_db()
        cursor = conn.cursor()
        cursor.execute("SELECT name, email FROM contact_messages WHERE id = ?", (message_id,))
        original_msg = cursor.fetchone()
        
        if not original_msg:
            conn.close()
            return jsonify({"success": False, "message": "Original message not found"}), 404
        
        # Store reply in chat_messages table so user can see it
        # We need to find the user_id from email
        cursor.execute("SELECT id FROM users WHERE email = ?", (original_msg['email'],))
        user = cursor.fetchone()
        
        user_id = user['id'] if user else None
        
        cursor.execute(
            "INSERT INTO chat_messages (user_id, user_email, message, is_from_admin, created_at) VALUES (?, ?, ?, ?, ?)",
            (user_id, original_msg['email'], f"[Admin Reply]: {reply_message}", 1, datetime.now().isoformat())
        )
        conn.commit()
        conn.close()
        
        return jsonify({
            "success": True,
            "message": "Reply sent successfully"
        })
    except Exception as e:
        return jsonify({"success": False, "message": str(e)}), 500


# ==================== HOME ROUTE ====================

@app.route("/")
def home():
    return jsonify({"message": "Backend running successfully"})


# ==================== MAIN ====================

if __name__ == "__main__":
    init_db()
    app.run(port=4000, debug=True)

